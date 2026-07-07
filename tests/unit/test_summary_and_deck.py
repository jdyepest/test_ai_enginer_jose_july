from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.integrations.llm_client import FakeLlmClient
from app.models.meeting_summary import MeetingSummary
from app.services.presentation_service import PresentationService


def test_fake_llm_response_matches_schema() -> None:
    summary = FakeLlmClient().summarize("BrightLane wants meeting intelligence.", "v1")

    assert isinstance(summary, MeetingSummary)
    assert len(summary.objectives) == 3
    assert len(summary.action_items) == 3
    assert summary.objectives[0].evidence[0].source_text


def test_powerpoint_has_required_slide_count_and_text(tmp_path: Path) -> None:
    summary = FakeLlmClient().summarize("BrightLane wants meeting intelligence.", "v1")
    deck_path = tmp_path / "deck.pptx"

    PresentationService().create_deck(summary, deck_path, "v1", "job-test")

    deck = Presentation(deck_path)
    all_text = "\n".join(
        shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert len(deck.slides) == 5
    assert "Executive summary" in all_text
    assert "Objectives" in all_text
    assert "Action Items" in all_text
    assert "Next Steps" in all_text

