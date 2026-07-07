from __future__ import annotations

import logging
from pathlib import Path

from app.models.meeting_summary import MeetingSummary
from app.utils.logging import log_event
from app.utils.timestamps import yyyymmdd

logger = logging.getLogger(__name__)


class PresentationService:
    def create_deck(
        self,
        summary: MeetingSummary,
        output_path: Path,
        prompt_version: str,
        job_id: str,
    ) -> Path:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        colors = {
            "ink": RGBColor(26, 32, 44),
            "muted": RGBColor(82, 92, 105),
            "line": RGBColor(20, 105, 126),
            "navy": RGBColor(18, 35, 46),
            "teal": RGBColor(38, 166, 154),
            "gold": RGBColor(232, 171, 74),
            "pale": RGBColor(232, 245, 247),
            "soft": RGBColor(247, 249, 251),
            "border": RGBColor(215, 224, 230),
            "white": RGBColor(255, 255, 255),
        }

        def style_shape(shape, fill_color, line_color=None, line_width: float = 0.75) -> None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.color.rgb = line_color or fill_color
            shape.line.width = Pt(line_width)

        def style_cell(cell, fill_color) -> None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

        def set_text(
            shape,
            text: str,
            size: int,
            color=None,
            *,
            bold: bool = False,
            align=None,
        ) -> None:
            frame = shape.text_frame
            frame.clear()
            frame.word_wrap = True
            paragraph = frame.paragraphs[0]
            paragraph.text = text
            paragraph.font.size = Pt(size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = color or colors["ink"]
            if align is not None:
                paragraph.alignment = align

        def add_textbox(
            slide,
            left: float,
            top: float,
            width: float,
            height: float,
            text: str,
            size: int = 14,
            color=None,
            *,
            bold: bool = False,
            align=None,
        ):
            box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            set_text(box, text, size, color, bold=bold, align=align)
            return box

        def add_section(slide, label: str, title: str) -> None:
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                slide_width,
                Inches(0.22),
            )
            style_shape(band, colors["teal"])
            add_textbox(slide, 0.55, 0.38, 4.5, 0.24, label.upper(), 8, colors["line"], bold=True)
            add_textbox(slide, 0.55, 0.65, 10.8, 0.48, title, 23, colors["ink"], bold=True)

        def add_footer(slide) -> None:
            add_textbox(
                slide,
                0.45,
                7.08,
                12.4,
                0.22,
                f"Generated {yyyymmdd()} | Prompt {prompt_version} | Human review required",
                8,
                colors["muted"],
                align=PP_ALIGN.RIGHT,
            )

        def add_bullet_list(
            slide,
            left: float,
            top: float,
            width: float,
            height: float,
            items: list[str],
            size: int = 13,
        ) -> None:
            box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            frame = box.text_frame
            frame.word_wrap = True
            frame.clear()
            for index, item in enumerate(items):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = item
                paragraph.level = 0
                paragraph.font.size = Pt(size)
                paragraph.font.color.rgb = colors["ink"]
                paragraph.space_after = Pt(7)

        def add_card(slide, left: float, top: float, width: float, height: float, fill=None):
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )
            style_shape(card, fill or colors["white"], colors["border"])
            return card

        def add_badge(slide, left: float, top: float, label: str, value: str, fill) -> None:
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(2.45),
                Inches(0.72),
            )
            style_shape(badge, fill, fill)
            add_textbox(
                slide,
                left + 0.15,
                top + 0.12,
                0.55,
                0.3,
                value,
                16,
                colors["white"],
                bold=True,
            )
            add_textbox(
                slide,
                left + 0.75,
                top + 0.16,
                1.45,
                0.3,
                label,
                9,
                colors["white"],
                bold=True,
            )

        def evidence_label(item) -> str:
            if not item.evidence:
                return "Evidence: not provided"
            first = item.evidence[0]
            prefix_parts = [part for part in [first.timestamp, first.speaker] if part]
            prefix = " | ".join(prefix_parts)
            source = truncate(first.source_text, 135)
            return f"{prefix}: {source}" if prefix else source

        def truncate(value: str, limit: int) -> str:
            clean = " ".join(value.split())
            if len(clean) <= limit:
                return clean
            return clean[: limit - 3].rstrip() + "..."

        blank = prs.slide_layouts[6]

        # Slide 1: executive summary.
        slide = prs.slides.add_slide(blank)
        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(3.05), slide_height)
        style_shape(sidebar, colors["navy"])
        add_textbox(
            slide,
            0.4,
            0.45,
            2.25,
            0.35,
            "MEETING INTELLIGENCE",
            9,
            colors["teal"],
            bold=True,
        )
        add_textbox(
            slide,
            0.38,
            1.02,
            2.35,
            1.2,
            summary.company or "Leadership Review",
            23,
            colors["white"],
            bold=True,
        )
        add_badge(slide, 0.38, 3.15, "objectives", "3", colors["line"])
        add_badge(slide, 0.38, 4.0, "action items", "3", colors["teal"])
        add_badge(slide, 0.38, 4.85, "review gate", "1", colors["gold"])
        add_textbox(slide, 3.55, 0.7, 8.9, 0.8, summary.meeting_title, 28, colors["ink"], bold=True)
        add_textbox(
            slide,
            3.58,
            1.55,
            2.4,
            0.25,
            "Executive summary",
            10,
            colors["line"],
            bold=True,
        )
        add_textbox(slide, 3.58, 1.9, 8.75, 2.95, summary.executive_summary, 17, colors["ink"])
        add_card(slide, 3.55, 5.25, 8.85, 0.76, colors["pale"])
        add_textbox(
            slide,
            3.78,
            5.43,
            8.25,
            0.28,
            "Output is ready for human review before distribution.",
            13,
            colors["line"],
            bold=True,
        )
        add_footer(slide)

        # Slide 2: objectives parsed into three readable cards.
        slide = prs.slides.add_slide(blank)
        add_section(slide, "Leadership Themes", "Three High-Level Objectives")
        card_positions = [(0.62, 1.55), (4.6, 1.55), (8.58, 1.55)]
        for index, objective in enumerate(summary.objectives, start=1):
            left, top = card_positions[index - 1]
            add_card(slide, left, top, 3.45, 4.75)
            add_textbox(
                slide,
                left + 0.22,
                top + 0.22,
                0.55,
                0.35,
                f"{index}",
                18,
                colors["teal"],
                bold=True,
            )
            add_textbox(
                slide,
                left + 0.22,
                top + 0.78,
                3.0,
                1.45,
                objective.objective,
                15,
                colors["ink"],
                bold=True,
            )
            add_textbox(
                slide,
                left + 0.22,
                top + 2.58,
                1.2,
                0.22,
                "Evidence",
                9,
                colors["line"],
                bold=True,
            )
            add_textbox(
                slide,
                left + 0.22,
                top + 2.9,
                2.98,
                1.2,
                evidence_label(objective),
                10,
                colors["muted"],
            )
        add_footer(slide)

        # Slide 3: action items as a structured table.
        slide = prs.slides.add_slide(blank)
        add_section(slide, "Execution", "Action Items by Priority")
        rows, cols = 4, 4
        table_shape = slide.shapes.add_table(
            rows,
            cols,
            Inches(0.55),
            Inches(1.45),
            Inches(12.2),
            Inches(4.85),
        )
        table = table_shape.table
        widths = [4.55, 1.25, 1.15, 5.25]
        headers = ["Action", "Owner", "Priority", "Business Rationale"]
        for index, width in enumerate(widths):
            table.columns[index].width = Inches(width)
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            style_cell(cell, colors["navy"])
            cell.text = header
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors["white"]
        for row, item in enumerate(summary.action_items, start=1):
            values = [
                item.action,
                item.owner or "TBD",
                item.priority.upper(),
                item.business_rationale,
            ]
            for col, value in enumerate(values):
                cell = table.cell(row, col)
                style_cell(cell, colors["soft"] if row % 2 else colors["white"])
                cell.text = truncate(value, 175 if col in {0, 3} else 24)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9 if col in {0, 3} else 10)
                paragraph.font.bold = col == 2
                paragraph.font.color.rgb = colors["ink"]
        add_footer(slide)

        # Slide 4: next steps as a simple timeline.
        slide = prs.slides.add_slide(blank)
        add_section(slide, "Review Path", "Next Steps and Human Review")
        timeline_y = 2.0
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2),
            Inches(timeline_y + 0.18),
            Inches(10.4),
            Inches(0.04),
        )
        style_shape(line, colors["border"])
        for index, step in enumerate(summary.next_steps, start=1):
            left = 0.72 + ((index - 1) * 4.05)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + 0.55),
                Inches(timeline_y),
                Inches(0.42),
                Inches(0.42),
            )
            style_shape(dot, colors["teal"])
            set_text(dot, str(index), 12, colors["white"], bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, left, 2.72, 3.35, 1.25, step.step, 14, colors["ink"], bold=True)
            add_textbox(
                slide,
                left,
                4.15,
                3.25,
                0.42,
                f"Owner: {step.owner or 'TBD'}",
                10,
                colors["muted"],
            )
            add_textbox(
                slide,
                left,
                4.58,
                3.25,
                0.42,
                f"Timeframe: {step.timeframe or 'TBD'}",
                10,
                colors["muted"],
            )
        add_card(slide, 0.75, 5.55, 11.6, 0.58, colors["pale"])
        add_textbox(
            slide,
            1.0,
            5.72,
            10.9,
            0.24,
            (
                "Final gate: reviewer confirms accuracy, evidence, audience, and "
                "distribution readiness."
            ),
            11,
            colors["line"],
            bold=True,
        )
        add_footer(slide)

        # Slide 5: audit details with evidence and risks separated.
        slide = prs.slides.add_slide(blank)
        add_section(slide, "Audit Trail", "Evidence and Risks")
        add_textbox(slide, 0.75, 1.35, 5.5, 0.28, "Source Evidence", 13, colors["line"], bold=True)
        evidence_items = []
        for objective in summary.objectives:
            first = objective.evidence[0]
            prefix = f"{first.timestamp} " if first.timestamp else ""
            speaker = f"{first.speaker}: " if first.speaker else ""
            evidence_items.append(f"{prefix}{speaker}{truncate(first.source_text, 165)}")
        add_bullet_list(slide, 0.75, 1.78, 5.65, 4.7, evidence_items, 10)
        add_textbox(
            slide,
            7.0,
            1.35,
            4.8,
            0.28,
            "Risks and Uncertainties",
            13,
            colors["gold"],
            bold=True,
        )
        add_bullet_list(
            slide,
            7.0,
            1.78,
            5.45,
            4.7,
            [truncate(risk, 155) for risk in summary.risks_and_uncertainties],
            10,
        )
        add_footer(slide)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        log_event(logger, "PRESENTATION_CREATED", job_id=job_id, deck_path=str(output_path))
        return output_path
